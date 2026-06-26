#!/usr/bin/env python3
"""Merge the 8 chapter decks into one master .pptx (image-aware slide copy +
sequential footer renumbering). Validates on save."""
import sys, re, copy
from io import BytesIO
sys.path.insert(0,"/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptxkit import validate_pptx

H=Path("/mnt/d/New folder/Antigravity-test/antigravity-skills/wedge_dossiers/book")
ORDER=["DataStaqAI-Book-Ch1-Executive-Thesis.pptx","DataStaqAI-Book-Ch2-AI-Transformation.pptx",
       "DataStaqAI-Book-Ch3-Business-Use-Cases.pptx","DataStaqAI-Book-Ch4a-Competitive-A.pptx",
       "DataStaqAI-Book-Ch4b-Competitive-B.pptx","DataStaqAI-Book-Ch5-Business-Model.pptx",
       "DataStaqAI-Book-Ch6-Path-to-5-10M.pptx","DataStaqAI-Book-Ch7-Investment-Thesis.pptx"]
EMBED=qn("r:embed"); LINK=qn("r:link")

def add_slide(master, src_slide):
    new=master.slides.add_slide(master.slide_layouts[6])
    for ph in list(new.shapes):
        ph._element.getparent().remove(ph._element)
    # map image relationships from source slide part -> new slide part
    rid_map={}
    for rid,rel in src_slide.part.rels.items():
        if rel.reltype==RT.IMAGE:
            ip=master.part.package.get_or_add_image_part(BytesIO(rel.target_part.blob))
            rid_map[rid]=new.part.relate_to(ip, RT.IMAGE)
    for shp in src_slide.shapes:
        el=copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in (EMBED,LINK):
                v=node.get(attr)
                if v in rid_map: node.set(attr,rid_map[v])
        new.shapes._spTree.append(el)
    if src_slide.has_notes_slide:
        new.notes_slide.notes_text_frame.text=src_slide.notes_slide.notes_text_frame.text
    return new

# base = chapter 1 (correct size + theme); then append the rest
master=Presentation(str(H/ORDER[0]))
for fn in ORDER[1:]:
    src=Presentation(str(H/fn))
    for sl in src.slides:
        add_slide(master, sl)
N=len(master.slides.__iter__.__self__._sldIdLst) if False else len(list(master.slides))

# renumber footers: find the "page / total" textbox and rewrite to "i/N"
pat=re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
for i,sl in enumerate(master.slides,1):
    for shp in sl.shapes:
        if shp.has_text_frame and pat.match(shp.text_frame.text.strip()):
            # rewrite first run, clear extras
            tf=shp.text_frame; p=tf.paragraphs[0]
            if p.runs:
                p.runs[0].text=f"{i}/{N}"
                for r in p.runs[1:]: r.text=""
            break

out=H/"DataStaqAI-The-Agentic-Enterprise.pptx"
master.save(str(out))
probs=validate_pptx(out)
print("MASTER slides:",N)
print("validation:", "OK" if not probs else probs)
print("OUTPUT:",out)
